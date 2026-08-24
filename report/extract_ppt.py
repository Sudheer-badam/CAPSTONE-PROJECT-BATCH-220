from pptx import Presentation
import sys

def extract_text(ppt_path):
    prs = Presentation(ppt_path)
    output = []
    for i, slide in enumerate(prs.slides):
        title = "Slide %d" % (i + 1)
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if shape != slide.shapes.title: # Don't duplicate title
                    slide_text.append(shape.text.strip())
        
        output.append(f"## {title}\n" + "\n".join(slide_text))
        
    with open('extracted_ppt_text.md', 'w', encoding='utf-8') as f:
        f.write("\n\n".join(output))

if __name__ == "__main__":
    extract_text("CAPSTONE PROJECT PPT-220-2 UPDATED.pptx")
